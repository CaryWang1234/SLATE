"""技能：文档安全扫描。

扫描文档文件（md/docx/pptx/xlsx/csv/pdf/txt），检测其中的不安全信息：
- 个人身份信息（PII）：身份证号、手机号、邮箱、地址
- 凭证与密钥：密码、API Key、Token
- 财务敏感数据：银行账号、薪资、金额
- 机密标记：「机密」「内部」「保密」等标记
- 内部信息：内网 URL、IP 地址、内部系统入口
- 敏感内容：健康/医疗记录
"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any

# ── 支持的文件格式 ────────────────────────────────
SCAN_EXTENSIONS = {
    ".md", ".txt", ".rst", ".log",          # 纯文本
    ".docx",                                 # Word
    ".pptx",                                 # PowerPoint
    ".xlsx", ".xlsm", ".csv",               # Excel / CSV
    ".pdf",                                 # PDF
}

EXCLUDE_DIRS = {
    "node_modules", ".git", "__pycache__", "venv", ".venv", "env",
    "dist", "build", ".next", ".nuxt", "target", "vendor",
}

MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 MB
MAX_FINDINGS_PER_FILE = 100


# ── 安全规则 ──────────────────────────────────────

class DocSecurityRule:
    """文档安全扫描规则。"""

    def __init__(self, pattern: re.Pattern, severity: str, category: str, description: str):
        self.pattern = pattern
        self.severity = severity   # critical / high / medium / low
        self.category = category
        self.description = description


# 个人身份信息 (PII)
PII_PATTERNS = [
    DocSecurityRule(
        re.compile(r"(?<!\d)[1-9]\d{5}(?:19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx](?!\d)"),
        "critical", "身份证号", "发现 18 位身份证号码"),
    DocSecurityRule(
        re.compile(r"(?<!\d)1[3-9]\d{9}(?!\d)"),
        "high", "手机号", "发现 11 位手机号码"),
    DocSecurityRule(
        re.compile(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}"),
        "medium", "邮箱地址", "发现电子邮箱地址"),
    DocSecurityRule(
        re.compile(r"(?:省|市|区|县|镇|乡|村|路|街|号|弄|栋|幢|单元|室)\s*.{0,5}(?:省|市|区|县|镇|乡|村|路|街|号|弄|栋|幢|单元|室)"),
        "medium", "详细地址", "发现详细居住地址"),
    DocSecurityRule(
        re.compile(r"(?:护照|护照号|Passport)\s*[:：]?\s*[A-Za-z0-9]{5,20}", re.I),
        "high", "护照号码", "发现护照号码"),
]

# 凭证与密钥
CREDENTIAL_PATTERNS = [
    DocSecurityRule(
        re.compile(r"(?:密码|口令|password|passwd|pwd)\s*[:：=]\s*\S{4,}", re.I),
        "critical", "硬编码密码", "文档中包含明文密码"),
    DocSecurityRule(
        re.compile(r"(?:api[_\-]?key|apikey|secret[_\-]?key|access[_\-]?token|auth[_\-]?token)\s*[:：=]\s*[\"']?\S{8,}[\"']?", re.I),
        "critical", "硬编码密钥", "文档中包含硬编码 API Key 或密钥"),
    DocSecurityRule(
        re.compile(r"(?:sk|pk|token|bearer)\s*[:：=]\s*[\"']?[A-Za-z0-9_\-]{16,}[\"']?", re.I),
        "high", "令牌/密钥片段", "疑似令牌或密钥片段"),
    DocSecurityRule(
        re.compile(r"-----BEGIN\s+(?:RSA|DSA|EC|OPENSSH|PGP)?\s*PRIVATE\s+KEY-----", re.I),
        "critical", "私钥内容", "文档中包含私钥内容"),
    DocSecurityRule(
        re.compile(r"(?:github|gitlab|bitbucket)[_\-]?(?:token|pat)\s*[:：=]\s*[\"']?\S{16,}[\"']?", re.I),
        "critical", "代码平台令牌", "文档中包含 Git 平台令牌"),
]

# 财务敏感数据
FINANCIAL_PATTERNS = [
    DocSecurityRule(
        re.compile(r"(?:银行|卡号|账户|account)\s*[:：]?\s*\d{10,25}", re.I),
        "critical", "银行账号", "发现银行账号信息"),
    DocSecurityRule(
        re.compile(r"(?:工资|薪资|薪酬|月薪|年薪|salary)\s*[:：]?\s*[¥$￥]?\s*[\d,]+(?:\.\d+)?(?:\s*(?:元|万|k|K))?", re.I),
        "high", "薪资信息", "发现薪资/薪酬数据"),
    DocSecurityRule(
        re.compile(r"(?:金额|总计|合计|费用|价格|报价|总价)\s*[:：]?\s*[¥$￥]\s*[\d,]+(?:\.\d+)?", re.I),
        "medium", "财务金额", "发现财务金额数据"),
    DocSecurityRule(
        re.compile(r"\b\d{4}[\s\-]?\d{4}[\s\-]?\d{4}[\s\-]?\d{4}\b"),
        "high", "银行卡号格式", "发现疑似银行卡号（16 位分组格式）"),
]

# 机密标记
CONFIDENTIAL_PATTERNS = [
    DocSecurityRule(
        re.compile(r"(?:绝密|机密|秘密|内部资料|保密|不得外传|禁止传播|confidential|secret|internal\s+only|do\s+not\s+distribute)", re.I),
        "high", "机密标记", "文档包含机密/内部标记"),
    DocSecurityRule(
        re.compile(r"(?:仅供.*?(?:内部|参考|审阅)|请勿.*?(?:外传|转发|复制)|此文档.*?敏感)", re.I),
        "medium", "分发限制标记", "文档包含分发限制说明"),
]

# 内部信息
INTERNAL_PATTERNS = [
    DocSecurityRule(
        re.compile(r"https?://(?:10\.\d{1,3}\.\d{1,3}\.\d{1,3}|172\.(?:1[6-9]|2\d|3[01])\.\d{1,3}\.\d{1,3}|192\.168\.\d{1,3}\.\d{1,3})(?::\d+)?(?:/\S*)?"),
        "high", "内网 URL", "发现内网 IP 地址或 URL"),
    DocSecurityRule(
        re.compile(r"(?:localhost|127\.0\.0\.1)\s*[:：]?\s*\d{2,5}"),
        "low", "本地服务地址", "发现本地服务地址"),
    DocSecurityRule(
        re.compile(r"(?:内网|VPN|代理|跳板机|堡垒机|admin[_\-]?panel|管理后台)\s*[:：]?\s*\S+", re.I),
        "medium", "内部系统信息", "发现内部系统入口信息"),
]

# 健康/医疗信息
HEALTH_PATTERNS = [
    DocSecurityRule(
        re.compile(r"(?:诊断|病历|处方|症状|用药|手术|住院|门诊|体检报告|medical\s*record|diagnosis|prescription)", re.I),
        "high", "医疗健康信息", "发现医疗健康记录"),
    DocSecurityRule(
        re.compile(r"(?:血型|HIV|乙肝|新冠|核酸|阳性|阴性|疫苗接种)", re.I),
        "medium", "健康检测信息", "发现健康检测相关信息"),
]

# 所有规则汇总
ALL_RULES = (
    PII_PATTERNS + CREDENTIAL_PATTERNS + FINANCIAL_PATTERNS +
    CONFIDENTIAL_PATTERNS + INTERNAL_PATTERNS + HEALTH_PATTERNS
)


# ── 文本提取 ──────────────────────────────────────

def _extract_text_txt(filepath: Path) -> str:
    """提取纯文本文件（.md/.txt/.rst/.log/.csv）。"""
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            return filepath.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return filepath.read_text(encoding="utf-8", errors="ignore")


def _extract_text_docx(filepath: Path) -> str:
    """提取 .docx 文本。"""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx 未安装")
    doc = Document(str(filepath))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_text_pptx(filepath: Path) -> str:
    """提取 .pptx 文本。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx 未安装")
    prs = Presentation(str(filepath))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if len(slide_texts) > 1:
            parts.append("\n".join(slide_texts))
    return "\n".join(parts)


def _extract_text_xlsx(filepath: Path) -> str:
    """提取 .xlsx 文本。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl 未安装")
    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        sheet_parts: list[str] = [f"[Sheet: {ws.title}]"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None]
            if cells:
                sheet_parts.append(" | ".join(cells))
        if len(sheet_parts) > 1:
            parts.append("\n".join(sheet_parts))
    wb.close()
    return "\n".join(parts)


def _extract_text_pdf(filepath: Path) -> str:
    """提取 .pdf 文本。"""
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber 未安装")
    parts: list[str] = []
    with pdfplumber.open(str(filepath)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_text(filepath: Path) -> str:
    """根据文件类型提取文本。"""
    ext = filepath.suffix.lower()
    if ext in (".md", ".txt", ".rst", ".log", ".csv"):
        return _extract_text_txt(filepath)
    if ext == ".docx":
        return _extract_text_docx(filepath)
    if ext == ".pptx":
        return _extract_text_pptx(filepath)
    if ext in (".xlsx", ".xlsm"):
        return _extract_text_xlsx(filepath)
    if ext == ".pdf":
        return _extract_text_pdf(filepath)
    return ""


# ── 扫描逻辑 ──────────────────────────────────────

def _should_skip_dir(dirpath: str) -> bool:
    """判断是否应跳过该目录。"""
    parts = Path(dirpath).parts
    return any(excluded in parts for excluded in EXCLUDE_DIRS)


def _scan_text(text: str, rules: list[DocSecurityRule], file_label: str,
               context_chars: int = 40) -> list[dict[str, Any]]:
    """扫描文本内容，返回发现的问题列表。"""
    findings: list[dict[str, Any]] = []
    lines = text.split("\n")

    for line_num, line in enumerate(lines, 1):
        if not line.strip():
            continue
        for rule in rules:
            for match in rule.pattern.finditer(line):
                # 提取上下文
                start = max(0, match.start() - context_chars)
                end = min(len(line), match.end() + context_chars)
                context = line[start:end].strip()
                if start > 0:
                    context = "..." + context
                if end < len(line):
                    context = context + "..."

                # 脱敏匹配内容（中间部分用 * 替换）
                matched = match.group()
                if len(matched) > 8:
                    masked = matched[:3] + "*" * (len(matched) - 6) + matched[-3:]
                else:
                    masked = matched[:2] + "***" + matched[-2:]

                findings.append({
                    "line": line_num,
                    "severity": rule.severity,
                    "category": rule.category,
                    "description": rule.description,
                    "matched": masked,
                    "context": context[:200],
                })
                if len(findings) >= MAX_FINDINGS_PER_FILE:
                    return findings
    return findings


def execute(
    directory: str = "",
    file_path: str = "",
    severity: str = "",
    category: str = "",
    max_files: int = 50,
    **_: Any,
) -> dict[str, Any]:
    """扫描文档中的不安全信息。

    参数：
    - directory: 扫描目录（递归扫描其中的文档文件）
    - file_path: 扫描单个文件（与 directory 二选一）
    - severity: 过滤最低严重级别 (critical/high/medium/low)
    - category: 过滤类别（如 "身份证号"、"硬编码密码"）
    - max_files: 最大扫描文件数（默认 50）
    """
    # 确定扫描目标
    if file_path:
        target = Path(file_path).resolve()
        if not target.is_file():
            return {"error": f"文件不存在: {file_path}"}
        if target.suffix.lower() not in SCAN_EXTENSIONS:
            return {"error": f"不支持的文件类型: {target.suffix}（支持 {', '.join(sorted(SCAN_EXTENSIONS))}）"}
        files_to_scan = [target]
        scan_dir = target.parent
    else:
        if not directory:
            directory = "."
        scan_dir = Path(directory).resolve()
        if not scan_dir.is_dir():
            return {"error": f"目录不存在: {directory}"}
        files_to_scan = []

    # 严重级别过滤
    severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    min_severity = severity_order.get(severity.lower(), 3) if severity else 3

    # 类别过滤
    if category:
        rules = [r for r in ALL_RULES if category.lower() in r.category.lower()]
        if not rules:
            available = sorted(set(r.category for r in ALL_RULES))
            return {"error": f"未知类别: {category}，可用: {', '.join(available)}"}
    else:
        rules = ALL_RULES

    # 收集文件（目录模式）
    skipped_count = 0
    if not file_path:
        for root, dirs, files in os.walk(scan_dir):
            if _should_skip_dir(root):
                dirs.clear()
                continue
            for fname in files:
                fpath = Path(root) / fname
                if fpath.suffix.lower() not in SCAN_EXTENSIONS:
                    continue
                try:
                    if fpath.stat().st_size > MAX_FILE_SIZE:
                        skipped_count += 1
                        continue
                except OSError:
                    continue
                files_to_scan.append(fpath)
                if len(files_to_scan) >= max_files:
                    break
            if len(files_to_scan) >= max_files:
                break

    # 逐文件扫描
    all_findings: list[dict[str, Any]] = []
    scanned_files: list[str] = []
    error_files: list[dict[str, str]] = []

    for fpath in files_to_scan:
        try:
            text = _extract_text(fpath)
        except RuntimeError as e:
            error_files.append({"file": str(fpath), "error": str(e)})
            continue
        except Exception as e:
            error_files.append({"file": str(fpath), "error": f"提取失败: {e}"})
            continue

        if not text.strip():
            continue

        scanned_files.append(str(fpath))
        file_findings = _scan_text(text, rules, str(fpath))

        for f in file_findings:
            f["file"] = str(fpath)
        all_findings.extend(file_findings)

    # 过滤严重级别
    filtered = [f for f in all_findings if severity_order.get(f["severity"], 3) <= min_severity]

    # 按严重级别排序
    filtered.sort(key=lambda x: severity_order.get(x["severity"], 3))

    # 统计
    severity_counts = {"critical": 0, "high": 0, "medium": 0, "low": 0}
    category_counts: dict[str, int] = {}
    for f in filtered:
        severity_counts[f["severity"]] = severity_counts.get(f["severity"], 0) + 1
        category_counts[f["category"]] = category_counts.get(f["category"], 0) + 1

    return {
        "scanned_files": len(scanned_files),
        "skipped_files": skipped_count,
        "error_files": error_files,
        "total_findings": len(filtered),
        "severity_summary": severity_counts,
        "category_summary": category_counts,
        "findings": filtered[:50],
        "truncated": len(filtered) > 50,
    }
