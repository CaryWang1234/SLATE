"""工具：PPT 制作 —— 根据主题、大纲与要点生成 .pptx 演示文稿。

参数:
- title      : 演示文稿标题（默认"未命名演示文稿"）
- subtitle   : 副标题（默认空）
- outline    : 大纲章节，逗号分隔，每个章节生成一页内容页（如"背景,方案,计划,总结"）
- slides     : 高级参数，JSON 数组精确控制每页内容，提供时忽略 outline，格式：
               [{"title": "章节标题", "points": ["要点一", "要点二"]}]
- theme      : 配色方案，内置 slate(黑金) / blue(深蓝) / green(墨绿) / wine(酒红) / gray(中性灰)，
               也可直接传 6 位十六进制色值（如 C9A96E）作为主色
- author     : 封面署名（默认空）
- output_dir : 输出目录（默认 SLATE 数据目录下的 outputs/）

实现：python-pptx 生成 16:9 文稿，封面页 + 内容页 + 结束页，主色用于装饰条与标题点缀。
"""

from __future__ import annotations

import json
import os
import re
import time
from datetime import date
from pathlib import Path
from typing import Any

DATA_DIR = Path(os.environ.get("SLATE_DATA_DIR", Path(__file__).resolve().parent.parent.parent / "data"))

# 内置配色：cover_bg 封面背景 / accent 主色 / cover_text 封面文字
THEMES: dict[str, dict[str, str]] = {
    "slate": {"cover_bg": "0A0A0A", "accent": "C9A96E", "cover_text": "E8DCC8"},
    "blue":  {"cover_bg": "0E2A47", "accent": "4C9BE8", "cover_text": "EAF2FB"},
    "green": {"cover_bg": "12332A", "accent": "4CAF7D", "cover_text": "E6F2EC"},
    "wine":  {"cover_bg": "3A1220", "accent": "C96A8B", "cover_text": "F5E6EC"},
    "gray":  {"cover_bg": "1F1F1F", "accent": "9E9E9E", "cover_text": "EDEDED"},
}
FONT_NAME = "Microsoft YaHei"


def execute(
    title: str = "",
    subtitle: str = "",
    outline: str = "",
    slides: Any = None,
    theme: str = "slate",
    author: str = "",
    output_dir: str = "",
    **_kw: Any,
) -> dict[str, Any]:
    """生成 .pptx 演示文稿，返回文件路径与统计信息。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return {"error": "缺少依赖 python-pptx，请先执行: pip install python-pptx"}

    doc_title = (title or "").strip() or "未命名演示文稿"
    theme_name, colors = _resolve_theme(theme)
    slide_list = _build_slides(slides, outline)

    accent = RGBColor.from_string(colors["accent"])
    cover_bg = RGBColor.from_string(colors["cover_bg"])
    cover_text = RGBColor.from_string(colors["cover_text"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 封面页 ──────────────────────────────
    cover = prs.slides.add_slide(blank)
    _fill_background(cover, cover_bg)
    _add_rect(cover, MSO_SHAPE.RECTANGLE, 0.9, 2.35, 1.2, 0.06, accent)
    _add_text(cover, 0.9, 2.6, 11.5, 1.6, [(doc_title, 44, True, cover_text)])
    if subtitle.strip():
        _add_text(cover, 0.9, 4.3, 11.5, 0.8, [(subtitle.strip(), 20, False, accent)])
    footer = " · ".join(x for x in [author.strip(), date.today().strftime("%Y-%m-%d")] if x)
    if footer:
        _add_text(cover, 0.9, 6.6, 11.5, 0.5, [(footer, 12, False, RGBColor.from_string("8A8A8A"))])

    # ── 内容页 ──────────────────────────────
    dark_text = RGBColor.from_string("1A1A1A")
    body_text = RGBColor.from_string("3C3C3C")
    for idx, item in enumerate(slide_list, 1):
        slide = prs.slides.add_slide(blank)
        _fill_background(slide, RGBColor.from_string("FFFFFF"))
        _add_rect(slide, MSO_SHAPE.RECTANGLE, 0.9, 0.5, 0.5, 0.07, accent)
        _add_text(slide, 0.9, 0.7, 11.5, 0.9, [(item["title"], 30, True, dark_text)])
        if item["points"]:
            lines = [(f"•  {p}", 18, False, body_text) for p in item["points"]]
            _add_text(slide, 1.0, 1.9, 11.3, 5.0, lines, space_after=14)
        _add_text(slide, 12.3, 6.9, 0.7, 0.4, [(str(idx), 11, False, RGBColor.from_string("AAAAAA"))],
                  align=PP_ALIGN.RIGHT)

    # ── 结束页 ──────────────────────────────
    ending = prs.slides.add_slide(blank)
    _fill_background(ending, cover_bg)
    _add_rect(ending, MSO_SHAPE.RECTANGLE, 6.07, 2.9, 1.2, 0.06, accent)
    _add_text(ending, 1.9, 3.2, 9.5, 1.4, [("感谢观看", 40, True, cover_text)], align=PP_ALIGN.CENTER)

    # ── 保存 ────────────────────────────────
    out_dir = Path(output_dir.strip()) if (output_dir or "").strip() else DATA_DIR / "outputs"
    try:
        out_dir.mkdir(parents=True, exist_ok=True)
        safe_name = re.sub(r'[\\/:*?"<>|\r\n]', "_", doc_title).strip()[:40] or "presentation"
        out_path = out_dir / f"{safe_name}_{time.strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(str(out_path))
    except OSError as e:
        return {"error": f"文件保存失败: {e}"}

    return {
        "file_path": str(out_path),
        "file_name": out_path.name,
        "size_bytes": out_path.stat().st_size,
        "slide_count": len(prs.slides),
        "theme": theme_name,
        "note": "演示文稿已生成，可用 PowerPoint / WPS 打开。",
    }


# ── 辅助函数 ──────────────────────────────────


def _resolve_theme(theme: str) -> tuple[str, dict[str, str]]:
    """解析配色：内置预设名或 6 位十六进制主色。"""
    t = (theme or "").strip().lower().lstrip("#")
    if re.fullmatch(r"[0-9a-fA-F]{6}", t):
        return f"custom({t.upper()})", {"cover_bg": "0A0A0A", "accent": t.upper(), "cover_text": "E8DCC8"}
    return t if t in THEMES else "slate", THEMES.get(t, THEMES["slate"])


def _build_slides(slides: Any, outline: str) -> list[dict[str, Any]]:
    """解析 slides JSON 或 outline 字符串为统一的页面列表。"""
    if slides:
        raw = slides
        if isinstance(raw, str):
            try:
                raw = json.loads(raw)
            except (ValueError, TypeError):
                raw = None
        if isinstance(raw, list):
            result = []
            for item in raw:
                if not isinstance(item, dict):
                    continue
                points = item.get("points") or []
                if isinstance(points, str):
                    points = [p.strip() for p in points.splitlines() if p.strip()]
                result.append({
                    "title": str(item.get("title") or "未命名章节"),
                    "points": [str(p) for p in points],
                })
            if result:
                return result

    chapters = [s.strip() for s in (outline or "").split(",") if s.strip()]
    return [{"title": ch, "points": []} for ch in chapters]


def _fill_background(slide: Any, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide: Any, shape_type: Any, left: float, top: float, width: float, height: float, color: Any) -> None:
    from pptx.util import Inches
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_text(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[tuple[str, int, bool, Any]],
    space_after: int = 0,
    align: Any = None,
) -> None:
    """添加文本框，lines 为 (文本, 字号, 是否加粗, 颜色) 列表。"""
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if space_after:
            para.space_after = Pt(space_after)
        if align is not None:
            para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT_NAME
        _set_east_asian_font(run)


def _set_east_asian_font(run: Any) -> None:
    """为中文字符设置东亚字体，避免回退到默认西文字体。"""
    from pptx.oxml.ns import qn
    from lxml import etree

    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", FONT_NAME)
